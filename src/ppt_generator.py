"""Generate EF-style PPT Memo using python-pptx."""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from .config import (
    EF_PURPLE_RGB, EF_ORANGE_RGB, EF_LIGHT_GRAY_RGB,
    EF_DARK_RGB, EF_WHITE_RGB, EF_BAR_BG_RGB,
)


def rgb(t):
    return RGBColor(*t)


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def txt(slide, left, top, w, h, text, size=18, color=EF_WHITE_RGB,
        bold=False, align=PP_ALIGN.LEFT, name="Arial"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = rgb(color)
    p.font.bold = bold
    p.font.name = name
    p.alignment = align
    return tf


def add_p(tf, text, size=16, color=EF_WHITE_RGB, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = rgb(color)
    p.font.bold = bold
    p.font.name = "Arial"
    p.space_before = Pt(10)
    return p


def bar(slide, left, top, w, score, label="", color=EF_ORANGE_RGB):
    """Score bar: label on left, bar in middle, score on right."""
    h = 0.28
    # Background bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = rgb(EF_BAR_BG_RGB)
    s.line.fill.background()
    # Filled portion
    fw = max(0.05, w * (score / 10))
    s2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(fw), Inches(h))
    s2.fill.solid()
    s2.fill.fore_color.rgb = rgb(color)
    s2.line.fill.background()
    # Label
    if label:
        txt(slide, left - 3.0, top - 0.04, 2.8, 0.4, label,
            size=16, color=EF_LIGHT_GRAY_RGB, align=PP_ALIGN.RIGHT)
    # Score
    txt(slide, left + w + 0.2, top - 0.04, 0.8, 0.4,
        f"{score}", size=16, color=EF_ORANGE_RGB, bold=True)


def header(slide, subtitle=""):
    """Add EF header to any slide."""
    txt(slide, 0.8, 0.4, 6, 0.5, "ENTREPRENEURS FIRST",
        size=20, color=EF_ORANGE_RGB, bold=True)
    if subtitle:
        txt(slide, 0.8, 1.0, 11, 0.7, subtitle,
            size=36, color=EF_WHITE_RGB, bold=True)


def generate_memo(
    candidate_name: str,
    alisa_eval: str,
    bob_eval: str,
    alisa_scores: dict,
    bob_scores: dict,
    round_table_result: dict,
    output_dir: str = "./output",
) -> str:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1: COVER (purple bg) =====
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s1, EF_PURPLE_RGB)

    txt(s1, 0.8, 0.4, 6, 0.5, "ENTREPRENEURS FIRST",
        size=20, color=EF_ORANGE_RGB, bold=True)

    txt(s1, 0.8, 2.0, 11, 1.5, candidate_name.upper(),
        size=56, color=EF_WHITE_RGB, bold=True)

    verdict = round_table_result.get("final_verdict", "Pending")
    edge = alisa_scores.get("EDGE_TYPE", "")
    txt(s1, 0.8, 3.8, 8, 0.7, f"VERDICT:  {verdict.upper()}",
        size=32, color=EF_ORANGE_RGB, bold=True)

    if edge:
        txt(s1, 0.8, 4.6, 8, 0.5, f"Edge:  {edge}",
            size=20, color=EF_LIGHT_GRAY_RGB)

    summary = alisa_scores.get("SUMMARY", "")
    if summary:
        txt(s1, 0.8, 5.5, 10, 1.0, summary,
            size=18, color=EF_LIGHT_GRAY_RGB)

    txt(s1, 0.8, 6.7, 11, 0.4,
        "Three AI minds. One decision. Zero bias.",
        size=14, color=EF_LIGHT_GRAY_RGB)

    # ===== SLIDE 2: FOUNDER EDGE — Alisa (dark bg) =====
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s2, EF_DARK_RGB)
    header(s2, "FOUNDER EDGE")

    txt(s2, 0.8, 1.8, 8, 0.5, "Evaluated by Alisa",
        size=16, color=EF_LIGHT_GRAY_RGB)

    dims = [
        ("Track Record", "TRACK_RECORD_SCORE"),
        ("Domain Expertise", "DOMAIN_EXPERTISE_SCORE"),
        ("Execution Signal", "EXECUTION_SIGNAL_SCORE"),
        ("Founder-Market Fit", "FOUNDER_MARKET_FIT_SCORE"),
    ]
    y = 2.6
    for label, key in dims:
        score = alisa_scores.get(key, 5)
        bar(s2, 3.8, y, 5.5, score, label=label)
        y += 0.75

    # Overall score
    overall = alisa_scores.get("OVERALL_SCORE", 5)
    txt(s2, 1.0, y + 0.3, 2.5, 0.5, "OVERALL",
        size=18, color=EF_LIGHT_GRAY_RGB, bold=True, align=PP_ALIGN.RIGHT)
    txt(s2, 3.8, y + 0.1, 2.5, 0.8, f"{overall}/10",
        size=40, color=EF_ORANGE_RGB, bold=True)

    verdict_a = alisa_scores.get("INITIAL_VERDICT", "N/A")
    txt(s2, 6.5, y + 0.25, 4, 0.6, f"→  {verdict_a}",
        size=24, color=EF_WHITE_RGB, bold=True)

    edge_text = alisa_scores.get("EDGE_TYPE", "")
    if edge_text:
        txt(s2, 0.8, 6.3, 11, 0.5, f"Edge Classification:  {edge_text}",
            size=18, color=EF_ORANGE_RGB, bold=True)

    # ===== SLIDE 3: TASTE & NETWORK — Bob (dark bg) =====
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s3, EF_DARK_RGB)
    header(s3, "TASTE & NETWORK")

    txt(s3, 0.8, 1.8, 8, 0.5, "Evaluated by Bob",
        size=16, color=EF_LIGHT_GRAY_RGB)

    dims_b = [
        ("Information Diet", "INFORMATION_DIET_SCORE"),
        ("Thought Leadership", "THOUGHT_LEADERSHIP_SCORE"),
        ("Network Quality", "NETWORK_QUALITY_SCORE"),
        ("Builder Signal", "BUILDER_SIGNAL_SCORE"),
    ]
    y = 2.6
    for label, key in dims_b:
        score = bob_scores.get(key, 5)
        bar(s3, 3.8, y, 5.5, score, label=label)
        y += 0.75

    overall_b = bob_scores.get("OVERALL_SCORE", 5)
    txt(s3, 1.0, y + 0.3, 2.5, 0.5, "OVERALL",
        size=18, color=EF_LIGHT_GRAY_RGB, bold=True, align=PP_ALIGN.RIGHT)
    txt(s3, 3.8, y + 0.1, 2.5, 0.8, f"{overall_b}/10",
        size=40, color=EF_ORANGE_RGB, bold=True)

    verdict_b = bob_scores.get("INITIAL_VERDICT", "N/A")
    txt(s3, 6.5, y + 0.25, 4, 0.6, f"→  {verdict_b}",
        size=24, color=EF_WHITE_RGB, bold=True)

    bob_summary = bob_scores.get("SUMMARY", "")
    if bob_summary:
        txt(s3, 0.8, 6.3, 11, 0.5, bob_summary[:100],
            size=16, color=EF_LIGHT_GRAY_RGB)

    # ===== SLIDE 4: ROUND TABLE (purple bg) =====
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s4, EF_PURPLE_RGB)
    header(s4, "ROUND TABLE")

    txt(s4, 0.8, 1.8, 11, 0.5,
        f"Debate  •  Alisa: {round_table_result['alisa_final_verdict']}  |  Bob: {round_table_result['bob_final_verdict']}",
        size=16, color=EF_LIGHT_GRAY_RGB)

    # Key turning points
    shifts = [r for r in round_table_result["rounds"] if r["shifted"]]
    y = 2.6
    txt(s4, 0.8, y, 4, 0.5, "KEY TURNING POINTS",
        size=18, color=EF_ORANGE_RGB, bold=True)
    y += 0.6

    for shift in shifts[:4]:
        short = shift["text"].split(".")[0] + "."
        if len(short) > 90:
            short = short[:87] + "..."
        agent_label = f"R{shift['round']}  {shift['agent']}"
        txt(s4, 0.8, y, 2.2, 0.5, agent_label,
            size=14, color=EF_ORANGE_RGB, bold=True)
        txt(s4, 3.1, y, 9.5, 0.5, short,
            size=14, color=EF_WHITE_RGB)
        y += 0.55

    # Final conviction
    y = max(y + 0.3, 5.0)
    txt(s4, 0.8, y, 5, 0.5, "FINAL CONVICTION",
        size=18, color=EF_ORANGE_RGB, bold=True)

    y += 0.55
    txt(s4, 0.8, y, 1.5, 0.4, "Alisa",
        size=16, color=EF_LIGHT_GRAY_RGB, bold=True)
    bar(s4, 2.5, y + 0.05, 4, round_table_result["alisa_final_conviction"], color=EF_PURPLE_RGB)
    txt(s4, 6.8, y, 4, 0.4,
        f"{round_table_result['alisa_final_conviction']}/10  →  {round_table_result['alisa_final_verdict']}",
        size=16, color=EF_WHITE_RGB, bold=True)

    y += 0.55
    txt(s4, 0.8, y, 1.5, 0.4, "Bob",
        size=16, color=EF_LIGHT_GRAY_RGB, bold=True)
    bar(s4, 2.5, y + 0.05, 4, round_table_result["bob_final_conviction"], color=EF_ORANGE_RGB)
    txt(s4, 6.8, y, 4, 0.4,
        f"{round_table_result['bob_final_conviction']}/10  →  {round_table_result['bob_final_verdict']}",
        size=16, color=EF_WHITE_RGB, bold=True)

    # ===== SLIDE 5: VERDICT (dark bg) =====
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s5, EF_DARK_RGB)
    header(s5)

    verdict = round_table_result.get("final_verdict", "Pending")
    txt(s5, 0.5, 2.0, 12.3, 1.5, verdict.upper(),
        size=72, color=EF_ORANGE_RGB, bold=True, align=PP_ALIGN.CENTER)

    consensus = "Unanimous" if round_table_result["consensus"] else "Split — Human Review Required"
    txt(s5, 0.5, 3.6, 12.3, 0.6, consensus,
        size=22, color=EF_LIGHT_GRAY_RGB, align=PP_ALIGN.CENTER)

    # Two columns
    tf1 = txt(s5, 1.5, 4.6, 4.5, 2.0, "ALISA — Founder Edge",
              size=18, color=EF_ORANGE_RGB, bold=True)
    add_p(tf1, f"Overall:  {alisa_scores.get('OVERALL_SCORE', '?')}/10",
          size=18, color=EF_WHITE_RGB, bold=True)
    add_p(tf1, f"Edge:  {alisa_scores.get('EDGE_TYPE', 'N/A')}",
          size=16, color=EF_LIGHT_GRAY_RGB)

    tf2 = txt(s5, 7.3, 4.6, 4.5, 2.0, "BOB — Taste & Network",
              size=18, color=EF_ORANGE_RGB, bold=True)
    add_p(tf2, f"Overall:  {bob_scores.get('OVERALL_SCORE', '?')}/10",
          size=18, color=EF_WHITE_RGB, bold=True)
    bob_sum = bob_scores.get("SUMMARY", "")
    if bob_sum:
        add_p(tf2, bob_sum[:70], size=16, color=EF_LIGHT_GRAY_RGB)

    txt(s5, 0.5, 6.7, 12.3, 0.4,
        "Three AI minds. One decision. Zero bias.  •  EF Agents",
        size=14, color=EF_LIGHT_GRAY_RGB, align=PP_ALIGN.CENTER)

    # Save
    os.makedirs(output_dir, exist_ok=True)
    safe = "".join(c if c.isalnum() or c in " -_" else "" for c in candidate_name)
    path = os.path.join(output_dir, f"{safe}_memo.pptx")
    prs.save(path)
    return path
